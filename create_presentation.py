#!/usr/bin/env python3
"""
Bitirme Projesi PowerPoint Sunumu Oluşturucu
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Kapak sayfası oluştur"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Başlık formatı
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def create_content_slide(prs, title, content_list):
    """İçerik slaytı oluştur"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    tf = body_shape.text_frame
    tf.clear()
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def create_two_column_slide(prs, title, left_content, right_content):
    """İki kolonlu slayt oluştur"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Başlık ekle
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Sol kolon
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(6)
    
    # Sağ kolon
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(6)
    
    return slide

def main():
    print("PowerPoint sunumu oluşturuluyor...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLAYT 1: Kapak
    create_title_slide(
        prs,
        "AKILLI SİTE YÖNETİM SİSTEMİ",
        "Entegre Web, Mobil ve Süper Admin Platformu\n\n[Öğrenci Adınız]\n[Danışman Adı]\nMayıs 2026"
    )
    
    # SLAYT 2: İçindekiler
    create_content_slide(prs, "İÇİNDEKİLER", [
        "1. Proje Tanımı ve Amacı",
        "2. Problem Tanımı",
        "3. Çözüm Yaklaşımı",
        "4. Sistem Mimarisi",
        "5. Temel Özellikler",
        "6. Kullanılan Teknolojiler",
        "7. Kullanıcı Rolleri ve Yetkiler",
        "8. Öne Çıkan Modüller",
        "9. Güvenlik ve Performans",
        "10. Demo ve Ekran Görüntüleri",
        "11. Sonuç ve Gelecek Planları"
    ])
    
    # SLAYT 3: Proje Tanımı
    create_content_slide(prs, "PROJE NEDİR?", [
        "📱 Akıllı Site Yönetim Sistemi - Apartman ve sitelerin tüm operasyonel süreçlerini dijitalleştiren kapsamlı platform",
        "",
        "🎯 3 Ana Platform:",
        "   • Mobil Uygulama (React Native)",
        "   • Web Yönetim Paneli (Next.js)",
        "   • Süper Admin Paneli",
        "",
        "🏢 Çoklu Site Yönetimi: Tek platformdan birden fazla site",
        "",
        "👥 Rol Bazlı Erişim: Admin, Sakin, Güvenlik, Temizlik, Süper Admin"
    ])
    
    # SLAYT 4: Proje Amacı
    create_two_column_slide(prs, "NEDEN BU PROJE?", 
        [
            "❌ PROBLEMLER:",
            "",
            "• Kağıt bazlı yönetim",
            "• İletişim kopuklukları",
            "• Aidat takip zorlukları",
            "• Paket yönetimi karmaşıklığı",
            "• Finansal raporlama eksikliği",
            "• Personel koordinasyon sorunları"
        ],
        [
            "✅ ÇÖZÜMLER:",
            "",
            "• Tamamen dijital süreç",
            "• Anlık bildirimler",
            "• Otomatik aidat hesaplama",
            "• QR kod paket sistemi",
            "• Detaylı finansal raporlar",
            "• Görev yönetimi sistemi"
        ]
    )
    
    # SLAYT 5: Sistem Mimarisi
    create_content_slide(prs, "SİSTEM MİMARİSİ", [
        "📱 KULLANICI KATMANI",
        "   Mobil App (React Native) | Web Panel (Next.js) | Super Admin",
        "",
        "🔌 API KATMANI",
        "   Spring Boot + Java | REST API | JWT Authentication",
        "",
        "🗄️ VERİTABANI KATMANI",
        "   MySQL Database | 40+ Tablo",
        "",
        "🔗 ENTEGRASYON SERVİSLERİ",
        "   Firebase (Push) | Gemini AI | Telegram"
    ])
    
    # SLAYT 6: Teknolojiler
    create_two_column_slide(prs, "KULLANILAN TEKNOLOJİLER",
        [
            "🔧 BACKEND:",
            "• Java 17 + Spring Boot 3.x",
            "• Spring Security + JWT",
            "• MySQL 8.0",
            "• Firebase Cloud Messaging",
            "• Google Gemini AI Vision",
            "",
            "💻 FRONTEND (WEB):",
            "• Next.js 14 (React)",
            "• Tailwind CSS + Shadcn UI",
            "• Recharts (Grafikler)",
            "• React Query"
        ],
        [
            "📱 MOBILE:",
            "• React Native + Expo",
            "• React Navigation",
            "• TypeScript",
            "• Push Notifications",
            "",
            "🚀 DEVOPS:",
            "• Docker",
            "• CI/CD Pipeline",
            "• Cloud Deployment Ready"
        ]
    )
    
    # SLAYT 7: Kullanıcı Rolleri
    create_content_slide(prs, "KULLANICI ROLLERİ VE YETKİLER", [
        "👑 SÜPER ADMIN: Tüm siteleri yönetme, sistem geneli raporlar",
        "",
        "👨‍💼 SİTE YÖNETİCİSİ: Sakin yönetimi, aidat takip, duyurular",
        "",
        "🏠 SAKİN: Aidat ödeme, arıza/talep, paket takibi, oylama",
        "",
        "🛡️ GÜVENLİK: Ziyaretçi onay, paket teslim, görev takibi",
        "",
        "🧹 TEMİZLİK: Görev listesi, tamamlanan işler, malzeme talepleri"
    ])
    
    # SLAYT 8: Temel Özellikler 1
    create_content_slide(prs, "TEMEL ÖZELLIKLER (1/2)", [
        "👥 SAKİN YÖNETİMİ: Malik/Kiracı ayrımı, çoklu daire, profil yönetimi",
        "",
        "💰 AİDAT YÖNETİMİ: Otomatik oluşturma, online ödeme, gecikme takibi",
        "",
        "📢 DUYURU SİSTEMİ: Öncelik bazlı, push notification, okunma takibi",
        "",
        "💬 MESAJLAŞMA: Sakin-Sakin, Sakin-Yönetim, blok bazlı mesajlar"
    ])
    
    # SLAYT 9: Temel Özellikler 2
    create_content_slide(prs, "TEMEL ÖZELLIKLER (2/2)", [
        "📦 PAKET YÖNETİMİ: QR kod, AI kargo tanıma, otomatik bildirim",
        "",
        "🔧 ARIZA/TALEP: Kategori bazlı, öncelik, durum takibi, fotoğraf",
        "",
        "🚪 ZİYARETÇİ YÖNETİMİ: QR kod giriş, zaman sınırlı, güvenlik onayı",
        "",
        "🗳️ E-OYLAMA: Dijital oylama, anonim/açık, sonuç grafikleri"
    ])
    
    # SLAYT 10: Özel Özellikler
    create_content_slide(prs, "ÖNE ÇIKAN YENİLİKLER", [
        "🤖 AI DESTEKLI KARGO TANIMA",
        "   Google Gemini Vision API | Otomatik bilgi çıkarma | %80+ doğruluk",
        "",
        "🏢 ÇOKLU SİTE YÖNETİMİ",
        "   Tek platform | Site bazlı izolasyon | Merkezi raporlama",
        "",
        "📊 FİNANSAL ANALİTİK",
        "   Gelir/Gider takibi | Dönemsel raporlar | Excel export",
        "",
        "✅ GÖREV YÖNETİMİ",
        "   Personel atama | Durum takibi | Performans raporları"
    ])
    
    # SLAYT 11: Güvenlik
    create_two_column_slide(prs, "GÜVENLİK VE PERFORMANS",
        [
            "🔒 GÜVENLİK:",
            "• JWT Token Authentication",
            "• Role-Based Access Control",
            "• Site bazlı veri izolasyonu",
            "• SQL Injection koruması",
            "• Input validation",
            "• Şifreli veri saklama",
            "• OTP doğrulama (Telegram)"
        ],
        [
            "⚡ PERFORMANS:",
            "• Optimize edilmiş sorgular",
            "• Database indexing",
            "• Lazy loading",
            "• Code splitting",
            "• Caching stratejileri",
            "• Real-time dashboard",
            "• <500ms response time"
        ]
    )
    
    # SLAYT 12: Veritabanı
    create_content_slide(prs, "VERİTABANI YAPISI", [
        "📊 40+ TABLO:",
        "   users, roles, sites, blocks, apartments, dues, payments,",
        "   announcements, messages, packages, tickets, visitor_requests,",
        "   tasks, voting, expenses, incomes, notifications...",
        "",
        "🔗 İLİŞKİLER:",
        "   • One-to-Many: Site → Blocks → Apartments",
        "   • Many-to-Many: Users ↔ Roles",
        "   • Soft Delete: is_deleted flag",
        "   • Audit: created_at, updated_at"
    ])
    
    # SLAYT 13: Mobil Uygulama
    create_content_slide(prs, "MOBİL UYGULAMA ÖZELLİKLERİ", [
        "📱 PLATFORM: iOS ve Android (React Native + Expo)",
        "",
        "✨ ÖZELLİKLER:",
        "   • Modern ve kullanıcı dostu arayüz",
        "   • Push notification",
        "   • Kamera entegrasyonu",
        "   • Interaktif grafikler",
        "   • Offline mode desteği",
        "   • Real-time senkronizasyon",
        "   • Çoklu dil desteği (TR/EN)"
    ])
    
    # SLAYT 14: Web Paneli
    create_two_column_slide(prs, "WEB YÖNETİM PANELİ",
        [
            "👨‍💼 ADMIN PANELİ:",
            "• Detaylı Dashboard",
            "• Sakin Yönetimi",
            "• Finansal Yönetim",
            "• Duyuru Yönetimi",
            "• Paket Yönetimi",
            "• Arıza/Talep Yönetimi",
            "• Raporlama ve Analitik",
            "• Site Ayarları"
        ],
        [
            "👑 SÜPER ADMIN:",
            "• Çoklu Site Yönetimi",
            "• Site Ekleme/Düzenleme",
            "• Sistem İstatistikleri",
            "• Toplu Mesajlaşma",
            "• Performans Analizi",
            "• Tüm Kullanıcılar",
            "• Sistem Ayarları"
        ]
    )
    
    # SLAYT 15: İstatistikler
    create_two_column_slide(prs, "PROJE İSTATİSTİKLERİ",
        [
            "📝 KOD İSTATİSTİKLERİ:",
            "• 40+ Veritabanı Tablosu",
            "• 50+ API Endpoint",
            "• 30+ Mobil Ekran",
            "• 25+ Web Sayfası",
            "• 5 Kullanıcı Rolü",
            "• 15+ Ana Modül",
            "• 2 Dil Desteği (TR/EN)"
        ],
        [
            "🧪 TEST VERİLERİ:",
            "• 400+ Test Kullanıcısı",
            "• 5 Farklı Site",
            "• 150+ Daire",
            "• Test Paketleri",
            "• Finansal Veriler",
            "• Görev Atamaları"
        ]
    )
    
    # SLAYT 16: Zorluklar ve Çözümler
    create_content_slide(prs, "KARŞILAŞILAN ZORLUKLAR VE ÇÖZÜMLER", [
        "1️⃣ Çoklu Site İzolasyonu → Site bazlı veri filtreleme",
        "",
        "2️⃣ Rol Bazlı Yetkilendirme → Spring Security + Custom Auth",
        "",
        "3️⃣ Real-time Bildirimler → Firebase Cloud Messaging",
        "",
        "4️⃣ AI Entegrasyonu → Google Gemini Vision API",
        "",
        "5️⃣ Performans Optimizasyonu → Indexing, caching, query optimization"
    ])
    
    # SLAYT 17: Test ve Doğrulama
    create_two_column_slide(prs, "TEST VE DOĞRULAMA",
        [
            "🧪 TEST SÜREÇLERİ:",
            "• Unit Testing",
            "• Integration Testing",
            "• API Testing (Postman)",
            "• UI/UX Testing",
            "• Security Testing",
            "• Performance Testing",
            "• User Acceptance Testing"
        ],
        [
            "✅ TEST SONUÇLARI:",
            "• %95+ API Success Rate",
            "• <500ms Response Time",
            "• Güvenlik Açığı: 0",
            "• Mobil Crash Rate: <1%",
            "• Tüm Kritik Özellikler OK"
        ]
    )
    
    # SLAYT 18: Gelecek Planları
    create_content_slide(prs, "GELECEK PLANLAR", [
        "📅 KISA VADELİ (3-6 ay):",
        "   iOS App Store | Chatbot | Gelişmiş analitik | SMS bildirimleri",
        "",
        "📅 ORTA VADELİ (6-12 ay):",
        "   Kurumsal paket | Multi-language | Tablet | 3. parti entegrasyonlar",
        "",
        "📅 UZUN VADELİ (1-2 yıl):",
        "   Uluslararası pazar | IoT entegrasyonu | Blockchain"
    ])
    
    # SLAYT 19: Sonuç
    create_two_column_slide(prs, "SONUÇ",
        [
            "✅ BAŞARILAR:",
            "• 3 platform tamamlandı",
            "• 15+ modül çalışıyor",
            "• Modern teknoloji yığını",
            "• Güvenli mimari",
            "• Kullanıcı dostu arayüz",
            "• AI entegrasyonu başarılı"
        ],
        [
            "📚 ÖĞRENİLENLER:",
            "• Full-stack development",
            "• Microservices mimarisi",
            "• Güvenlik best practices",
            "• Cross-platform dev",
            "• AI/ML entegrasyonu",
            "• Proje yönetimi"
        ]
    )
    
    # SLAYT 20: Teşekkür
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Başlık
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "TEŞEKKÜRLER"
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Alt metin
    content_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2))
    content_frame = content_box.text_frame
    content_frame.text = "Dinlediğiniz için teşekkür ederim!\n\nSorularınızı yanıtlamaktan mutluluk duyarım."
    content_frame.paragraphs[0].font.size = Pt(24)
    content_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Dosyayı kaydet
    filename = "Akilli_Site_Yonetim_Sistemi_Sunum.pptx"
    prs.save(filename)
    print(f"\n✅ Sunum başarıyla oluşturuldu: {filename}")
    print(f"📊 Toplam {len(prs.slides)} slayt oluşturuldu")
    print("\n💡 İpucu: Sunumu açıp ekran görüntüleri ve görseller ekleyebilirsiniz!")

if __name__ == "__main__":
    try:
        main()
    except ImportError:
        print("❌ Hata: python-pptx kütüphanesi bulunamadı!")
        print("\n📦 Kurulum için:")
        print("   pip install python-pptx")
    except Exception as e:
        print(f"❌ Hata oluştu: {e}")
